from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(r"C:\Users\kseno\Desktop\ДИПЛОМ\Презентация_защита_ВК_банк.pptx")


BG = RGBColor(11, 18, 32)
PANEL = RGBColor(16, 26, 42)
PANEL_2 = RGBColor(19, 31, 48)
CARD = RGBColor(29, 48, 73)
BORDER = RGBColor(35, 54, 82)
BORDER_2 = RGBColor(46, 72, 106)
ACCENT = RGBColor(114, 160, 255)
TEXT = RGBColor(244, 247, 255)
MUTED = RGBColor(170, 188, 214)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_box(slide, left, top, width, height, fill=PANEL, line=BORDER, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Segoe UI",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title_block(slide, kicker, title, subtitle=""):
    add_text(slide, Inches(0.45), Inches(0.2), Inches(3.2), Inches(0.22), kicker, 11, ACCENT, True)
    add_text(slide, Inches(0.45), Inches(0.45), Inches(8.2), Inches(0.6), title, 24, TEXT, True)
    if subtitle:
        add_text(slide, Inches(0.45), Inches(0.93), Inches(8.2), Inches(0.28), subtitle, 11, MUTED)


def style_center_text(shape, text, size=14, color=TEXT, bold=True):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Segoe UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_pill(slide, left, top, width, height, text):
    shape = add_box(slide, left, top, width, height, PANEL_2, BORDER_2, True)
    style_center_text(shape, text, 12, TEXT, False)
    return shape


def add_mock_phone(slide, left, top, width, height, label):
    add_box(slide, left, top, width, height, PANEL, BORDER, True)
    add_text(slide, left + Inches(0.1), top + Inches(0.06), width - Inches(0.2), Inches(0.16), label, 10, ACCENT, True)
    inner_left = left + Inches(0.1)
    inner_top = top + Inches(0.18)
    inner_width = width - Inches(0.2)
    for i in range(3):
        add_box(slide, inner_left, inner_top + Inches(0.26 * i), inner_width, Inches(0.18), CARD, BORDER_2, True)
    add_box(slide, inner_left, inner_top + Inches(0.82), inner_width, Inches(0.58), CARD, BORDER_2, True)
    start = inner_top + Inches(1.5)
    for i in range(4):
        add_box(slide, inner_left, start + Inches(0.22 * i), inner_width, Inches(0.14), PANEL_2, BORDER_2, True)


def add_mock_desktop(slide, left, top, width, height, label):
    add_box(slide, left, top, width, height, PANEL, BORDER, True)
    add_text(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.18), label, 10, ACCENT, True)
    add_box(slide, left + Inches(0.15), top + Inches(0.28), width - Inches(0.3), Inches(0.18), CARD, BORDER_2, True)
    add_box(slide, left + Inches(0.15), top + Inches(0.56), Inches(1.15), height - Inches(0.72), PANEL_2, BORDER_2, True)
    add_box(slide, left + Inches(1.42), top + Inches(0.56), width - Inches(1.57), Inches(0.7), PANEL_2, BORDER_2, True)
    add_box(slide, left + Inches(1.42), top + Inches(1.38), width - Inches(1.57), height - Inches(1.54), PANEL_2, BORDER_2, True)


def add_connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = ACCENT
    line.line.width = Pt(2)
    line.line.end_arrowhead = True
    return line


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(0.45), Inches(0.34), Inches(3.8), Inches(0.2), "ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА", 11, ACCENT, True)
add_text(
    slide,
    Inches(0.45),
    Inches(0.75),
    Inches(7.2),
    Inches(1.2),
    "Разработка информационной системы обслуживания клиентов банка на платформе ВКонтакте",
    26,
    TEXT,
    True,
)
add_text(
    slide,
    Inches(0.45),
    Inches(2.0),
    Inches(4.0),
    Inches(0.8),
    "Специальность: 09.02.07\nИнформационные системы и программирование\nВыполнил: Киреев Даниил Витальевич",
    15,
    MUTED,
)
add_text(slide, Inches(0.45), Inches(3.0), Inches(2.6), Inches(0.3), "Калининград, 2026", 13, MUTED)
add_mock_phone(slide, Inches(7.45), Inches(1.55), Inches(1.8), Inches(3.75), "VK Mini App")
add_mock_desktop(slide, Inches(9.45), Inches(1.55), Inches(3.1), Inches(3.75), "Административная панель")
add_pill(slide, Inches(7.45), Inches(5.45), Inches(1.8), Inches(0.32), "Клиентский контур")
add_pill(slide, Inches(9.45), Inches(5.45), Inches(3.1), Inches(0.32), "Внутренний контур сотрудников")

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_block(slide, "ПОСТАНОВКА ЗАДАЧИ", "Актуальность, цель и задачи")
add_box(slide, Inches(0.45), Inches(1.55), Inches(6.2), Inches(3.7))
add_text(slide, Inches(0.7), Inches(1.8), Inches(2.0), Inches(0.3), "Актуальность", 19, TEXT, True)
add_text(
    slide,
    Inches(0.7),
    Inches(2.2),
    Inches(5.5),
    Inches(1.5),
    "• Рост цифровых банковских сервисов\n• Удобство взаимодействия внутри экосистемы VK\n• Необходимость объединить клиентский и административный контуры",
    18,
    MUTED,
)
add_box(slide, Inches(6.95), Inches(1.55), Inches(5.9), Inches(3.7))
add_text(slide, Inches(7.2), Inches(1.8), Inches(1.4), Inches(0.3), "Цель", 19, TEXT, True)
add_text(
    slide,
    Inches(7.2),
    Inches(2.18),
    Inches(5.2),
    Inches(0.75),
    "Разработать информационную систему обслуживания клиентов банка на платформе ВКонтакте.",
    18,
    MUTED,
)
add_text(slide, Inches(7.2), Inches(3.0), Inches(1.6), Inches(0.3), "Задачи", 19, TEXT, True)
add_text(
    slide,
    Inches(7.2),
    Inches(3.38),
    Inches(5.2),
    Inches(1.35),
    "• Спроектировать архитектуру системы\n• Реализовать mini app, backend и БД\n• Создать административную панель\n• Реализовать поддержку и тестирование",
    17,
    MUTED,
)

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_block(slide, "АРХИТЕКТУРА", "Состав и взаимодействие компонентов")
blocks = [
    (Inches(0.7), Inches(1.7), Inches(2.0), Inches(0.8), "VK Mini App\nКлиентский интерфейс"),
    (Inches(3.2), Inches(1.7), Inches(2.0), Inches(0.8), "Backend\nFastAPI + бизнес-логика"),
    (Inches(5.7), Inches(1.7), Inches(2.0), Inches(0.8), "PostgreSQL\nХранение данных"),
    (Inches(9.2), Inches(1.15), Inches(2.4), Inches(0.8), "Admin Panel\nСотрудники банка"),
    (Inches(9.2), Inches(2.2), Inches(2.4), Inches(0.8), "VK-бот\nДополнительный канал"),
    (Inches(9.2), Inches(3.25), Inches(2.4), Inches(0.8), "AI-помощник\nПервичная поддержка"),
]
for left, top, width, height, text in blocks:
    shape = add_box(slide, left, top, width, height)
    style_center_text(shape, text, 15, TEXT, True)
add_connector(slide, Inches(2.7), Inches(2.1), Inches(3.2), Inches(2.1))
add_connector(slide, Inches(5.2), Inches(2.1), Inches(5.7), Inches(2.1))
add_connector(slide, Inches(7.7), Inches(2.1), Inches(9.2), Inches(1.55))
add_connector(slide, Inches(7.7), Inches(2.1), Inches(9.2), Inches(2.6))
add_connector(slide, Inches(7.7), Inches(2.1), Inches(9.2), Inches(3.65))
add_text(
    slide,
    Inches(0.7),
    Inches(4.65),
    Inches(7.8),
    Inches(0.9),
    "Единая серверная часть обслуживает пользовательское mini app, административную панель и бот. Это позволяет централизовать бизнес-логику, проверки и доступ к данным.",
    17,
    MUTED,
)

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_block(slide, "ПОЛЬЗОВАТЕЛЬСКОЕ ПРИЛОЖЕНИЕ", "Клиентское VK Mini App", "Основные экраны и пользовательские сценарии.")
add_mock_phone(slide, Inches(0.7), Inches(1.7), Inches(2.1), Inches(4.1), "Главный экран")
add_mock_phone(slide, Inches(3.05), Inches(1.7), Inches(2.1), Inches(4.1), "Счета и карты")
add_mock_phone(slide, Inches(5.4), Inches(1.7), Inches(2.1), Inches(4.1), "Переводы и заявки")
add_box(slide, Inches(8.0), Inches(1.7), Inches(4.75), Inches(4.1))
add_text(slide, Inches(8.25), Inches(1.95), Inches(2.1), Inches(0.3), "Что умеет клиент", 19, TEXT, True)
add_text(
    slide,
    Inches(8.25),
    Inches(2.35),
    Inches(4.0),
    Inches(2.3),
    "• Просмотр счетов и карт\n• История операций\n• Переводы между своими счетами\n• Перевод по VK ID\n• Межбанковский перевод\n• Подача заявок на продукты\n• Поддержка и уведомления",
    16,
    MUTED,
)
add_pill(slide, Inches(8.25), Inches(5.05), Inches(1.55), Inches(0.34), "PIN-защита")
add_pill(slide, Inches(9.95), Inches(5.05), Inches(1.9), Inches(0.34), "VK-интеграция")

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_block(slide, "ФУНКЦИОНАЛ", "Ключевые возможности системы")
features = [
    ("Счета", "Создание и просмотр счетов"),
    ("Карты", "Реквизиты и управление картами"),
    ("Операции", "История и детализация операций"),
    ("Переводы", "Свои счета, VK ID, межбанк"),
    ("Заявки", "Подача и отслеживание статусов"),
    ("Поддержка", "Чат и сервисные запросы"),
    ("Уведомления", "Лента уведомлений и статусы"),
    ("Безопасность", "PIN, роли, аудит, CSRF"),
]
start_x = Inches(0.55)
start_y = Inches(1.65)
box_w = Inches(3.0)
box_h = Inches(1.25)
gap_x = Inches(0.17)
gap_y = Inches(0.18)
for i, (name, desc) in enumerate(features):
    col = i % 4
    row = i // 4
    left = start_x + col * (box_w + gap_x)
    top = start_y + row * (box_h + gap_y)
    add_box(slide, left, top, box_w, box_h)
    add_text(slide, left + Inches(0.18), top + Inches(0.18), box_w - Inches(0.36), Inches(0.25), name, 18, TEXT, True)
    add_text(slide, left + Inches(0.18), top + Inches(0.52), box_w - Inches(0.36), Inches(0.5), desc, 13, MUTED)
add_text(
    slide,
    Inches(0.62),
    Inches(4.72),
    Inches(12.0),
    Inches(0.45),
    "Система покрывает полный клиентский контур: от просмотра продуктов до переводов, заявок, поддержки и настроек безопасности.",
    16,
    MUTED,
)

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_block(slide, "АДМИНИСТРАТИВНАЯ ПАНЕЛЬ", "Внутренний контур сотрудников банка")
add_mock_desktop(slide, Inches(0.55), Inches(1.7), Inches(7.0), Inches(3.9), "Сводка системы, метрики и графики")
add_mock_desktop(slide, Inches(7.8), Inches(1.7), Inches(4.95), Inches(3.9), "Аудит и управление сотрудниками")
add_pill(slide, Inches(0.7), Inches(5.8), Inches(1.25), Inches(0.34), "operator")
add_pill(slide, Inches(2.05), Inches(5.8), Inches(1.15), Inches(0.34), "admin")
add_pill(slide, Inches(3.3), Inches(5.8), Inches(1.6), Inches(0.34), "superadmin")
add_text(
    slide,
    Inches(5.15),
    Inches(5.75),
    Inches(7.1),
    Inches(0.55),
    "Роли, cookie-сессии, CSRF-защита, аудит действий, обработка заявок и сервисных запросов.",
    15,
    MUTED,
)

# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_block(slide, "ТЕХНИЧЕСКАЯ РЕАЛИЗАЦИЯ", "Выбранный стек и ключевые технологии")
tech = [
    ("Mini App", "React\nVite\nVK Bridge"),
    ("Admin Panel", "React\nRecharts"),
    ("Backend", "Python\nFastAPI\nSQLAlchemy"),
    ("База данных", "PostgreSQL"),
    ("Бот", "Python\nvk_api"),
    ("Инфраструктура", "Docker Compose\nNginx"),
]
positions = [
    (Inches(0.55), Inches(1.8)),
    (Inches(4.1), Inches(1.8)),
    (Inches(7.65), Inches(1.8)),
    (Inches(0.55), Inches(4.0)),
    (Inches(4.1), Inches(4.0)),
    (Inches(7.65), Inches(4.0)),
]
for (title, body), (left, top) in zip(tech, positions):
    add_box(slide, left, top, Inches(2.95), Inches(1.55))
    add_text(slide, left + Inches(0.2), top + Inches(0.16), Inches(2.5), Inches(0.25), title, 18, TEXT, True)
    add_text(slide, left + Inches(0.2), top + Inches(0.5), Inches(2.5), Inches(0.8), body, 14, MUTED)

# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title_block(slide, "РЕЗУЛЬТАТЫ", "Итоги работы и практическая значимость")
add_box(slide, Inches(0.55), Inches(1.7), Inches(12.25), Inches(3.6))
add_text(slide, Inches(0.85), Inches(2.0), Inches(3.0), Inches(0.3), "Что реализовано", 21, TEXT, True)
add_text(
    slide,
    Inches(0.85),
    Inches(2.42),
    Inches(4.9),
    Inches(2.0),
    "• VK Mini App для клиентов\n• Backend с банковской бизнес-логикой\n• PostgreSQL как единое хранилище\n• Административная панель сотрудников\n• VK-бот и AI-модуль поддержки",
    17,
    MUTED,
)
add_text(slide, Inches(6.2), Inches(2.0), Inches(4.0), Inches(0.3), "Практическая значимость", 21, TEXT, True)
add_text(
    slide,
    Inches(6.2),
    Inches(2.42),
    Inches(5.8),
    Inches(2.0),
    "• Проект показывает построение многокомпонентной системы\n• Объединяет клиентский и внутренний контуры\n• Демонстрирует архитектурную целостность\n• Может быть расширен до более зрелой версии",
    17,
    MUTED,
)
for i, label in enumerate(["Клиентский контур", "Административный контур", "Единый backend", "Безопасность"]):
    add_pill(slide, Inches(0.85 + i * 2.95), Inches(5.55), Inches(2.55), Inches(0.36), label)
add_text(slide, Inches(4.4), Inches(6.35), Inches(4.6), Inches(0.4), "Спасибо за внимание", 26, TEXT, True, PP_ALIGN.CENTER)

OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT_PATH))
print(OUTPUT_PATH)
